"""
GenAI Security Gateway - Dosya Analiz Kontrolcüsü (File Controller)

Sisteme yüklenen dosyaların (PDF, TXT, DOCX vb.) içeriklerini okur ve 
içlerinde güvenlik zafiyeti veya hassas veri (PII) olup olmadığını kontrol eder.
Özellikle büyük dosyaların (örneğin 100 sayfalık bir PDF) güvenli bir şekilde 
ayrıştırılıp (parsing) güvenlik katmanlarından (Layer 1, 2, 3) geçirilmesini sağlar.
"""

from fastapi import APIRouter, UploadFile, File, HTTPException
import logging
import io

logger = logging.getLogger(__name__)

router = APIRouter()

@router.post("/extract-document", tags=["Document Processing"])
async def extract_document(file: UploadFile = File(...)):
    """
    Extracts text from PDF, DOCX, and PPTX files.
    """
    filename = file.filename.lower()
    
    try:
        content = await file.read()
        file_obj = io.BytesIO(content)
        
        extracted_text = ""
        
        if filename.endswith(".pdf"):
            import pypdf
            reader = pypdf.PdfReader(file_obj)
            for page in reader.pages:
                extracted_text += page.extract_text() + "\n"
                
        elif filename.endswith(".docx"):
            import docx
            doc = docx.Document(file_obj)
            for para in doc.paragraphs:
                extracted_text += para.text + "\n"
                
        elif filename.endswith(".pptx"):
            from pptx import Presentation
            prs = Presentation(file_obj)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        extracted_text += shape.text + "\n"
        elif filename.endswith((".png", ".jpg", ".jpeg", ".webp", ".bmp")):
            from PIL import Image
            from app.services.llm_proxy import LLMProxy
            
            image = Image.open(file_obj)
            client = LLMProxy.get_client("gemini-2.5-flash")
            if not client:
                raise HTTPException(status_code=503, detail="Yapay zeka servisi kullanılamıyor (Kota dolu olabilir).")
            
            prompt = "Lütfen bu görseldeki tüm metinleri tam olarak ve eksiksiz çıkar. Sadece metni döndür, ekstra yorum ekleme."
            response = await client.aio.models.generate_content(
                model='gemini-2.5-flash',
                contents=[image, prompt]
            )
            
            extracted_text = response.text.strip() if response.text else "[Görselden metin çıkarılamadı]"
            
        else:
            raise HTTPException(status_code=400, detail="Unsupported file format.")
            
        return {"filename": file.filename, "extracted_text": extracted_text.strip()}
        
    except Exception as e:
        logger.error(f"Error extracting document {filename}: {e}", exc_info=True)
        raise HTTPException(status_code=500, detail=f"File processing error: {str(e)}")
